from pathlib import Path
from typing import Optional



def truncate_slides(file_path: str, max_pages: int, output_path: Optional[str] = None) -> str:
    """
    Truncate a PDF or PPTX file to a given page/slide count, keeping only the first max_pages pages.
    
    Args:
        file_path: Path to a PDF or PPTX file.
        max_pages: Maximum number of pages/slides to keep.
        output_path: Output file path. If None, a new file name is generated:
            {stem}_truncated_to_{max_pages}{suffix}
        
    Returns:
        Path to the output file.
        
    Raises:
        FileNotFoundError: If the file does not exist.
        ValueError: If the file format is not supported or max_pages <= 0.
    """
    file_path = Path(file_path)
    
    if not file_path.exists():
        raise FileNotFoundError(f"File does not exist: {file_path}")
    
    if max_pages <= 0:
        raise ValueError(f"max_pages must be > 0, current value: {max_pages}")
    
    file_ext = file_path.suffix.lower()
    
    if output_path is None:
        # Generate new file name: {stem}_truncated_to_{max_pages}{suffix}
        stem = file_path.stem
        suffix = file_path.suffix
        new_filename = f"{stem}_truncated_to_{max_pages}{suffix}"
        output_path = str(file_path.parent / new_filename)
    else:
        output_path = str(Path(output_path))
    
    if file_ext == '.pdf':
        return _truncate_pdf(file_path, max_pages, output_path)
    elif file_ext == '.pptx':
        return _truncate_pptx(file_path, max_pages, output_path)
    else:
        raise ValueError(f"Unsupported file format: {file_ext}. Only .pdf and .pptx are supported.")


def _truncate_pdf(file_path: Path, max_pages: int, output_path: str) -> str:
    """Truncate a PDF file to max_pages pages."""
    try:
        import PyPDF2
    except ImportError:
        try:
            import pypdf as PyPDF2
        except ImportError:
            raise ImportError(
                "PyPDF2 or pypdf is required to process PDF files. "
                "Install with: pip install pypdf2  (or)  pip install pypdf"
            )
    
    with open(file_path, 'rb') as file:
        pdf_reader = PyPDF2.PdfReader(file)
        total_pages = len(pdf_reader.pages)
        
        if total_pages <= max_pages:
            # No truncation needed; return original file path.
            return str(file_path)
        
        # Create a new PDF writer.
        pdf_writer = PyPDF2.PdfWriter()
        
        # Add only the first max_pages pages.
        for page_num in range(max_pages):
            pdf_writer.add_page(pdf_reader.pages[page_num])
        
        # Write to output file.
        with open(output_path, 'wb') as output_file:
            pdf_writer.write(output_file)
        
        return output_path


def _truncate_pptx(file_path: Path, max_pages: int, output_path: str) -> str:
    """Truncate a PPTX file to max_pages slides."""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError(
            "python-pptx is required to process PPTX files. "
            "Install with: pip install python-pptx"
        )
    
    prs = Presentation(file_path)
    total_slides = len(prs.slides)
    
    if total_slides <= max_pages:
        # No truncation needed; return original file path.
        return str(file_path)
    
    # Delete slides beyond max_pages (delete from the end).
    slide_ids = prs.slides._sldIdLst
    while len(slide_ids) > max_pages:
        # Get the rId of the last slide.
        last_slide_id = slide_ids[-1]
        rId = last_slide_id.rId
        # Drop relationship.
        prs.part.drop_rel(rId)
        # Remove from list.
        slide_ids.remove(last_slide_id)
    
    # Save to output file.
    prs.save(output_path)
    
    return output_path



if __name__ == "__main__":
    file_path = '/path/to/slides.pdf'
    max_pages = 10
    output_path = truncate_slides(file_path, max_pages)
    print(output_path)
